import streamlit as st
import requests
from pptx import Presentation
import tempfile

def main():
    st.title("PPT from GitHub Demo with Navigation")

    ppt_github_url = "https://github.com/ShayneGeo/ppt/raw/refs/heads/main/NIROPS_DatasetDescription.pptx"

    # Initialize session state variables
    if "slides_text" not in st.session_state:
        st.session_state.slides_text = []
        st.session_state.current_slide = 0

    if st.button("Load PPT from GitHub"):
        response = requests.get(ppt_github_url)
        if response.status_code == 200:
            # Write to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name

            # Parse the PPTX
            pres = Presentation(tmp_path)
            slides_text = []
            for slide in pres.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                slides_text.append("\n".join(slide_content))

            # Store text in session state
            st.session_state.slides_text = slides_text
            st.session_state.current_slide = 0
        else:
            st.error("Failed to download PPT from GitHub")

    # Only show navigation if we have slides
    if st.session_state.slides_text:
        total_slides = len(st.session_state.slides_text)

        col1, col2 = st.columns(2)
        with col1:
            # Disable 'Previous' if on the first slide
            if st.button("Previous", disabled=(st.session_state.current_slide == 0)):
                st.session_state.current_slide -= 1

        with col2:
            # Disable 'Next' if on the last slide
            if st.button("Next", disabled=(st.session_state.current_slide == total_slides - 1)):
                st.session_state.current_slide += 1

        # Show which slide we're on
        st.write(f"Slide {st.session_state.current_slide + 1} of {total_slides}")

        # Display the current slide's text
        st.write(st.session_state.slides_text[st.session_state.current_slide])

if __name__ == "__main__":
    main()
